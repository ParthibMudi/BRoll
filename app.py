import os
import json
import uuid
import subprocess
import tempfile
import shutil
from flask import Flask, request, jsonify, send_file, send_from_directory
from flask_cors import CORS
from groq import Groq
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont
import math
import cv2
import numpy as np
from sklearn.cluster import KMeans

app = Flask(__name__)
CORS(app)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
INPUT_DIR = os.path.join(BASE_DIR, "input")
OUTPUT_DIR = os.path.join(BASE_DIR, "output")

os.makedirs(INPUT_DIR, exist_ok=True)
os.makedirs(OUTPUT_DIR, exist_ok=True)

# ─── GROQ B-ROLL PLANNING ────────────────────────────────────────────────────

def analyze_transcript_with_groq(transcript_data: dict, api_key: str) -> dict:
    client = Groq(api_key=api_key)
    
    segments_text = "\n".join([
        f"[{seg['id']}] ({seg['start']}s - {seg['end']}s): {seg['text']}"
        for seg in transcript_data["segments"]
    ])
    
    prompt = f"""You are a video production assistant. Analyze this transcript and identify 1-3 segments that would benefit from visual B-roll explanation.

TRANSCRIPT:
{segments_text}

RULES:
- Only select segments with: key insights, lists/enumerations, comparisons, processes, decision points, problem/solution, or concepts easier to show visually
- SKIP: greetings, transitions, fillers, simple statements, repetitions
- For each selected segment, choose a template: "four_point_list" OR "decision_tree" OR "problem_solution"
- four_point_list: for lists, reasons, benefits, steps (provide exactly 4 points)
- decision_tree: for yes/no decisions, comparisons, conditions (provide one root question with yes/no branches)
- problem_solution: for problem + solution pairs (provide problem and solution text)

Respond ONLY with valid JSON in this exact format:
{{
  "video_id": "{transcript_data.get('video_id', 'video_001')}",
  "broll_items": [
    {{
      "id": "broll_001",
      "start": 12.0,
      "end": 18.0,
      "source_segment_ids": ["seg_002"],
      "selected_text": "the exact segment text",
      "template_type": "four_point_list",
      "title": "Slide Title Here",
      "points": ["Point 1", "Point 2", "Point 3", "Point 4"],
      "reason_for_selection": "Why this segment needs visual support"
    }}
  ],
  "skipped_segments": [
    {{
      "segment_id": "seg_001",
      "text": "segment text",
      "reason": "Why this was skipped"
    }}
  ]
}}

For decision_tree template use:
"nodes": [{{"label": "Question?", "yes": "Yes outcome", "no": "No outcome"}}]
instead of "points"

For problem_solution template use:
"problem": "Problem text", "solution": "Solution text"
instead of "points"

Be selective. Only pick the strongest 1-3 moments. Respond ONLY with JSON, no markdown."""

    response = client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.3,
        max_tokens=2000
    )
    
    raw = response.choices[0].message.content.strip()
    # Strip markdown fences if present
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
    raw = raw.strip()
    
    return json.loads(raw)


# ─── PPT GENERATION ──────────────────────────────────────────────────────────

BRAND_COLORS = {
    "bg":      RGBColor(0x0D, 0x0D, 0x1A),
    "accent":  RGBColor(0x6C, 0x63, 0xFF),
    "accent2": RGBColor(0x00, 0xD9, 0xA6),
    "white":   RGBColor(0xFF, 0xFF, 0xFF),
    "light":   RGBColor(0xB0, 0xAE, 0xD8),
    "card":    RGBColor(0x1A, 0x1A, 0x2E),
}
def extract_palette_from_video(video_path, num_keyframes=5, n_colors=5):
    # Extract frames with ffmpeg to memory or temp files
    frames = []
    cap = cv2.VideoCapture(video_path)
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    for i in range(num_keyframes):
        sec = (i + 1) * (get_video_duration(video_path) / (num_keyframes + 1))
        cap.set(cv2.CAP_PROP_POS_MSEC, sec * 1000)
        ret, frame = cap.read()
        if ret:
            frames.append(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB))
    cap.release()
    
    # Combine all pixels (downsampled for speed)
    all_pixels = np.vstack([cv2.resize(f, (160, 90)).reshape(-1, 3) for f in frames])
    
    # K-Means clustering
    kmeans = KMeans(n_clusters=n_colors, random_state=42, n_init=10).fit(all_pixels)
    colors = kmeans.cluster_centers_.astype(int)
    
    # Count cluster frequency
    labels, counts = np.unique(kmeans.labels_, return_counts=True)
    ordered = [colors[idx] for idx in labels[np.argsort(-counts)]]
    
    # Convert to hex or RGBColor (skip too gray colors if desired)
    palette = []
    for c in ordered:
        if not (abs(int(c[0])-int(c[1]))<20 and abs(int(c[1])-int(c[2]))<20):  # not gray
            palette.append(RGBColor(*c))
    return palette[:3]   # keep top 3 vibrant colors

def set_bg(slide, prs):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BRAND_COLORS["bg"]

def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color or BRAND_COLORS["white"]
    return txBox

def add_rounded_rect(slide, left, top, width, height, color, radius=0.08):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def generate_four_point_list_ppt(item: dict, output_path: str):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, prs)
    
    # Top accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.07))
    bar.fill.solid(); bar.fill.fore_color.rgb = BRAND_COLORS["accent"]
    bar.line.fill.background()
    
    # Title
    add_text_box(slide, item["title"],
                 0.4, 0.18, 9.2, 0.8,
                 font_size=28, bold=True,
                 color=BRAND_COLORS["white"], align=PP_ALIGN.LEFT)
    
    # Subtitle line
    sub = slide.shapes.add_shape(1, Inches(0.4), Inches(0.95), Inches(1.2), Inches(0.04))
    sub.fill.solid(); sub.fill.fore_color.rgb = BRAND_COLORS["accent2"]
    sub.line.fill.background()
    
    points = item.get("points", ["", "", "", ""])[:4]
    emojis = ["01", "02", "03", "04"]
    positions = [(0.35, 1.15), (5.2, 1.15), (0.35, 3.0), (5.2, 3.0)]
    
    for i, (pt, pos) in enumerate(zip(points, positions)):
        lx, ly = pos
        # Card bg
        card = slide.shapes.add_shape(1,
            Inches(lx), Inches(ly), Inches(4.5), Inches(1.65))
        card.fill.solid(); card.fill.fore_color.rgb = BRAND_COLORS["card"]
        card.line.color.rgb = BRAND_COLORS["accent"]
        card.line.width = Pt(1)
        
        # Number badge
        badge = slide.shapes.add_shape(1,
            Inches(lx+0.12), Inches(ly+0.15), Inches(0.5), Inches(0.5))
        badge.fill.solid(); badge.fill.fore_color.rgb = BRAND_COLORS["accent"]
        badge.line.fill.background()
        
        add_text_box(slide, emojis[i],
                     lx+0.12, ly+0.15, 0.5, 0.5,
                     font_size=13, bold=True,
                     color=BRAND_COLORS["white"], align=PP_ALIGN.CENTER)
        
        add_text_box(slide, pt,
                     lx+0.75, ly+0.2, 3.6, 1.3,
                     font_size=17, bold=False,
                     color=BRAND_COLORS["white"])
    
    # Footer
    add_text_box(slide, "AI-Generated Visual Explainer",
                 0.3, 5.25, 9.4, 0.3,
                 font_size=9, color=BRAND_COLORS["light"], align=PP_ALIGN.RIGHT)
    
    prs.save(output_path)


def generate_decision_tree_ppt(item: dict, output_path: str):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, prs)
    
    # Accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.07))
    bar.fill.solid(); bar.fill.fore_color.rgb = BRAND_COLORS["accent2"]
    bar.line.fill.background()
    
    # Title
    add_text_box(slide, item["title"],
                 0.4, 0.18, 9.2, 0.8,
                 font_size=28, bold=True,
                 color=BRAND_COLORS["white"])
    
    node = item.get("nodes", [{"label": "Question?", "yes": "Yes", "no": "No"}])[0]
    
    # Root box
    root = slide.shapes.add_shape(1, Inches(3.2), Inches(1.2), Inches(3.6), Inches(1.0))
    root.fill.solid(); root.fill.fore_color.rgb = BRAND_COLORS["accent"]
    root.line.fill.background()
    add_text_box(slide, node["label"],
                 3.2, 1.2, 3.6, 1.0,
                 font_size=16, bold=True, align=PP_ALIGN.CENTER)
    
    # YES branch (left)
    yes_box = slide.shapes.add_shape(1, Inches(0.8), Inches(3.3), Inches(3.5), Inches(0.9))
    yes_box.fill.solid(); yes_box.fill.fore_color.rgb = RGBColor(0xE5, 0x73, 0x73)
    yes_box.line.fill.background()
    add_text_box(slide, "YES  →  " + node["yes"],
                 0.8, 3.3, 3.5, 0.9,
                 font_size=15, bold=False, align=PP_ALIGN.CENTER)
    
    # YES label
    add_text_box(slide, "YES", 1.8, 2.55, 1.0, 0.4,
                 font_size=12, bold=True,
                 color=RGBColor(0xE5, 0x73, 0x73), align=PP_ALIGN.CENTER)
    
    # NO branch (right)
    no_box = slide.shapes.add_shape(1, Inches(5.7), Inches(3.3), Inches(3.5), Inches(0.9))
    no_box.fill.solid(); no_box.fill.fore_color.rgb = BRAND_COLORS["accent2"]
    no_box.line.fill.background()
    add_text_box(slide, "NO  →  " + node["no"],
                 5.7, 3.3, 3.5, 0.9,
                 font_size=15, bold=False, align=PP_ALIGN.CENTER)
    
    # NO label
    add_text_box(slide, "NO", 7.2, 2.55, 1.0, 0.4,
                 font_size=12, bold=True,
                 color=BRAND_COLORS["accent2"], align=PP_ALIGN.CENTER)
    
    # Connector lines using thin rectangles
    # Vertical down from root
    v1 = slide.shapes.add_shape(1, Inches(4.98), Inches(2.2), Inches(0.04), Inches(0.65))
    v1.fill.solid(); v1.fill.fore_color.rgb = BRAND_COLORS["light"]; v1.line.fill.background()
    
    # Horizontal split
    h = slide.shapes.add_shape(1, Inches(2.55), Inches(2.82), Inches(4.9), Inches(0.04))
    h.fill.solid(); h.fill.fore_color.rgb = BRAND_COLORS["light"]; h.line.fill.background()
    
    # Down to YES
    v2 = slide.shapes.add_shape(1, Inches(2.55), Inches(2.82), Inches(0.04), Inches(0.5))
    v2.fill.solid(); v2.fill.fore_color.rgb = BRAND_COLORS["light"]; v2.line.fill.background()
    
    # Down to NO
    v3 = slide.shapes.add_shape(1, Inches(7.41), Inches(2.82), Inches(0.04), Inches(0.5))
    v3.fill.solid(); v3.fill.fore_color.rgb = BRAND_COLORS["light"]; v3.line.fill.background()
    
    # Footer
    add_text_box(slide, "AI-Generated Visual Explainer",
                 0.3, 5.25, 9.4, 0.3,
                 font_size=9, color=BRAND_COLORS["light"], align=PP_ALIGN.RIGHT)
    
    prs.save(output_path)


def generate_problem_solution_ppt(item: dict, output_path: str):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, prs)
    
    # Accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.07))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(0xFF, 0x6B, 0x6B)
    bar.line.fill.background()
    
    # Title
    add_text_box(slide, item["title"],
                 0.4, 0.18, 9.2, 0.8,
                 font_size=28, bold=True, color=BRAND_COLORS["white"])
    
    problem = item.get("problem", "Problem not specified")
    solution = item.get("solution", "Solution not specified")
    
    # Problem card
    p_card = slide.shapes.add_shape(1, Inches(0.4), Inches(1.2), Inches(4.2), Inches(3.5))
    p_card.fill.solid(); p_card.fill.fore_color.rgb = RGBColor(0x2A, 0x10, 0x10)
    p_card.line.color.rgb = RGBColor(0xFF, 0x6B, 0x6B); p_card.line.width = Pt(1.5)
    
    add_text_box(slide, "⚠ PROBLEM", 0.55, 1.3, 4.0, 0.5,
                 font_size=13, bold=True, color=RGBColor(0xFF, 0x6B, 0x6B))
    add_text_box(slide, problem, 0.55, 1.85, 3.9, 2.5,
                 font_size=18, color=BRAND_COLORS["white"])
    
    # Arrow
    add_text_box(slide, "→", 4.75, 2.5, 0.6, 0.8,
                 font_size=36, bold=True,
                 color=BRAND_COLORS["accent2"], align=PP_ALIGN.CENTER)
    
    # Solution card
    s_card = slide.shapes.add_shape(1, Inches(5.4), Inches(1.2), Inches(4.2), Inches(3.5))
    s_card.fill.solid(); s_card.fill.fore_color.rgb = RGBColor(0x0A, 0x2A, 0x1A)
    s_card.line.color.rgb = BRAND_COLORS["accent2"]; s_card.line.width = Pt(1.5)
    
    add_text_box(slide, "✓ SOLUTION", 5.55, 1.3, 4.0, 0.5,
                 font_size=13, bold=True, color=BRAND_COLORS["accent2"])
    add_text_box(slide, solution, 5.55, 1.85, 3.9, 2.5,
                 font_size=18, color=BRAND_COLORS["white"])
    
    # Footer
    add_text_box(slide, "AI-Generated Visual Explainer",
                 0.3, 5.25, 9.4, 0.3,
                 font_size=9, color=BRAND_COLORS["light"], align=PP_ALIGN.RIGHT)
    
    prs.save(output_path)


def generate_ppt(item: dict, output_path: str):
    t = item.get("template_type", "four_point_list")
    if t == "four_point_list":
        generate_four_point_list_ppt(item, output_path)
    elif t == "decision_tree":
        generate_decision_tree_ppt(item, output_path)
    elif t == "problem_solution":
        generate_problem_solution_ppt(item, output_path)
    else:
        generate_four_point_list_ppt(item, output_path)


# ─── SLIDE → PNG → MP4 ───────────────────────────────────────────────────────

_W, _H = 1920, 1080
_BG_BASE = (8, 8, 20)
_WHITE   = (255, 255, 255)
_OFF_WHITE = (220, 218, 255)
_MUTED   = (120, 118, 160)


def _font(size: int, bold: bool = False):
    """Return a PIL TrueType font, falling back gracefully."""
    try:
        candidates = (["arialbd.ttf", "calibrib.ttf", "segoeui.ttf"] if bold
                      else ["arial.ttf", "calibri.ttf", "segoeui.ttf", "verdana.ttf"])
        for name in candidates:
            for root in [r"C:\Windows\Fonts",
                         os.path.expanduser("~/AppData/Local/Microsoft/Windows/Fonts")]:
                p = os.path.join(root, name)
                if os.path.exists(p):
                    return ImageFont.truetype(p, size)
    except Exception:
        pass
    return ImageFont.load_default()


def _draw_text_wrapped(draw, text, x, y, max_w, font, fill, line_spacing=1.45):
    words = text.split()
    lines, cur = [], ""
    for w in words:
        test = (cur + " " + w).strip()
        bb = draw.textbbox((0, 0), test, font=font)
        if bb[2] - bb[0] > max_w and cur:
            lines.append(cur); cur = w
        else:
            cur = test
    if cur:
        lines.append(cur)
    for line in lines:
        draw.text((x, y), line, font=font, fill=fill)
        bb = draw.textbbox((0, 0), line, font=font)
        y += int((bb[3] - bb[1]) * line_spacing)
    return y


def _blend(c1, c2, t):
    """Linear blend between two RGB tuples."""
    return tuple(int(c1[i] + (c2[i] - c1[i]) * t) for i in range(3))


def _darken(c, factor=0.35):
    return tuple(int(v * factor) for v in c)


def _lighten(c, factor=0.25):
    return tuple(min(255, int(v + (255 - v) * factor)) for v in c)


def _draw_gradient_bg(img, top_color, bot_color):
    """Draw a vertical gradient background."""
    arr = np.array(img)
    h = arr.shape[0]
    for y in range(h):
        t = y / h
        r = int(top_color[0] + (bot_color[0] - top_color[0]) * t)
        g = int(top_color[1] + (bot_color[1] - top_color[1]) * t)
        b = int(top_color[2] + (bot_color[2] - top_color[2]) * t)
        arr[y, :] = [r, g, b]
    return Image.fromarray(arr.astype(np.uint8))


def _draw_card(draw, x1, y1, x2, y2, fill, border_color, border_w=3, shadow_offset=6):
    """Draw a card with a subtle drop-shadow effect."""
    # Shadow
    shadow = tuple(max(0, v - 30) for v in fill)
    draw.rectangle([x1 + shadow_offset, y1 + shadow_offset,
                    x2 + shadow_offset, y2 + shadow_offset], fill=shadow)
    # Card body
    draw.rectangle([x1, y1, x2, y2], fill=fill, outline=border_color, width=border_w)


def _draw_pill_badge(draw, cx, cy, r, fill, text, font):
    """Draw a circular badge with number."""
    draw.ellipse([cx - r, cy - r, cx + r, cy + r], fill=fill)
    bb = draw.textbbox((0, 0), text, font=font)
    tw, th = bb[2] - bb[0], bb[3] - bb[1]
    draw.text((cx - tw // 2, cy - th // 2 - 2), text, font=font, fill=_WHITE)


def _draw_connector(draw, x1, y1, x2, y2, color, width=4):
    draw.line([(x1, y1), (x2, y2)], fill=color, width=width)


def render_broll_item_to_png(item: dict, png_path: str, palette: list = None):
    """
    Beautiful Pillow renderer with video-extracted color palette.
    palette: list of up to 3 RGB tuples from the source video.
    """
    # ── Resolve accent colors from video palette or defaults ──────────────────
    default_acc1 = (108, 99, 255)
    default_acc2 = (0, 217, 166)
    default_acc3 = (255, 107, 107)

    acc1 = tuple(int(v) for v in palette[0]) if palette and len(palette) > 0 else default_acc1
    acc2 = tuple(int(v) for v in palette[1]) if palette and len(palette) > 1 else default_acc2
    acc3 = tuple(int(v) for v in palette[2]) if palette and len(palette) > 2 else default_acc3

    # Ensure enough contrast — boost saturation if the color is too dark
    def _boost(c):
        mx = max(c)
        if mx < 80:
            scale = 180 / max(mx, 1)
            return tuple(min(255, int(v * scale)) for v in c)
        return c

    acc1, acc2, acc3 = _boost(acc1), _boost(acc2), _boost(acc3)

    bg_top = _darken(acc1, 0.07)
    bg_bot = (6, 6, 14)
    card_bg = _blend(_BG_BASE, acc1, 0.08)

    t     = item.get("template_type", "four_point_list")
    title = item.get("title", "")

    # ── Base gradient background ──────────────────────────────────────────────
    img = Image.new("RGB", (_W, _H), bg_top)
    img = _draw_gradient_bg(img, bg_top, bg_bot)
    draw = ImageDraw.Draw(img)

    # Subtle diagonal grid pattern overlay
    for gx in range(0, _W, 60):
        draw.line([(gx, 0), (gx, _H)], fill=(*acc1[:3], ), width=1)
        if gx < _H:
            draw.line([(0, gx), (_W, gx)], fill=(*acc1[:3], ), width=1)
    # Re-apply gradient on top of grid (blend)
    overlay = Image.new("RGBA", (_W, _H), (0, 0, 0, 0))
    od = ImageDraw.Draw(overlay)
    for y in range(_H):
        alpha = int(200 * (1 - y / _H * 0.4))
        od.line([(0, y), (_W, y)], fill=(*bg_top, alpha))
    img = img.convert("RGBA")
    img = Image.alpha_composite(img, overlay).convert("RGB")
    draw = ImageDraw.Draw(img)

    # ── Top accent bar (thick gradient strip) ─────────────────────────────────
    for bx in range(_W):
        t_frac = bx / _W
        bar_c = _blend(acc1, acc2, t_frac)
        draw.line([(bx, 0), (bx, 7)], fill=bar_c)
    # Glow below bar
    for bx in range(_W):
        t_frac = bx / _W
        glow_c = _blend(acc1, acc2, t_frac)
        a = 60
        draw.line([(bx, 8), (bx, 18)], fill=(*glow_c, ))

    # ── Title area ────────────────────────────────────────────────────────────
    title_font  = _font(72, bold=True)
    sub_font    = _font(32)
    badge_font  = _font(40, bold=True)
    body_font   = _font(38)
    label_font  = _font(30, bold=True)
    small_font  = _font(26)

    # Title shadow
    draw.text((78, 42), title, font=title_font, fill=_darken(acc1, 0.5))
    draw.text((76, 40), title, font=title_font, fill=_WHITE)

    # Accent underline
    bb = draw.textbbox((76, 40), title, font=title_font)
    ul_w = min(bb[2] - bb[0], 300)
    for ux in range(ul_w):
        draw.line([(76 + ux, bb[3] + 6), (76 + ux, bb[3] + 12)],
                  fill=_blend(acc1, acc2, ux / ul_w))

    # Template label tag
    tag_labels = {"four_point_list": "4-POINT LIST", "decision_tree": "DECISION TREE",
                  "problem_solution": "PROBLEM / SOLUTION"}
    tag_text = tag_labels.get(t, t.upper())
    tag_col = acc1 if t == "four_point_list" else (acc2 if t == "decision_tree" else acc3)
    tag_bg = _darken(tag_col, 0.3)
    tag_x = _W - 320
    draw.rectangle([tag_x, 30, _W - 30, 80], fill=tag_bg, outline=tag_col, width=2)
    draw.text((tag_x + 14, 38), tag_text, font=small_font, fill=tag_col)

    # ── Template bodies ───────────────────────────────────────────────────────
    if t == "four_point_list":
        points   = item.get("points", [])[:4]
        nums     = ["01", "02", "03", "04"]
        cols     = [acc1, acc2, acc3, _blend(acc1, acc2, 0.5)]
        pos      = [(60, 190), (1000, 190), (60, 610), (1000, 610)]

        for i, (pt, (px, py)) in enumerate(zip(points, pos)):
            c = cols[i % len(cols)]
            c_dark = _darken(c, 0.18)
            # Card
            _draw_card(draw, px, py, px + 870, py + 370, card_bg, c, border_w=3)
            # Left accent strip
            for sx in range(8):
                t_frac = sx / 8
                draw.line([(px + sx, py), (px + sx, py + 370)],
                          fill=_blend(c, _darken(c, 0.4), t_frac))
            # Number badge
            _draw_pill_badge(draw, px + 70, py + 60, 44, c, nums[i], badge_font)
            # Point text
            _draw_text_wrapped(draw, pt, px + 130, py + 28, 710, body_font, _OFF_WHITE)

    elif t == "decision_tree":
        node = item.get("nodes", [{"label": "?", "yes": "Yes", "no": "No"}])[0]
        cx = _W // 2

        # Root question box
        _draw_card(draw, cx - 400, 175, cx + 400, 390, _darken(acc1, 0.25), acc1, border_w=4)
        for lx in range(8):
            draw.line([(cx - 400 + lx, 175), (cx - 400 + lx, 390)],
                      fill=_blend(acc1, _darken(acc1, 0.3), lx / 8))
        _draw_text_wrapped(draw, node["label"], cx - 370, 205, 740, _font(44, bold=True), _WHITE)

        # Connector lines
        _draw_connector(draw, cx, 390, cx, 490, _MUTED)
        _draw_connector(draw, 450, 490, _W - 450, 490, _MUTED)
        _draw_connector(draw, 450, 490, 450, 570, _MUTED)
        _draw_connector(draw, _W - 450, 490, _W - 450, 570, _MUTED)

        # YES label
        draw.text((310, 498), "YES ✓", font=label_font, fill=acc2)
        # NO label
        draw.text((_W - 510, 498), "NO ✗", font=label_font, fill=acc3)

        # YES outcome box
        _draw_card(draw, 70, 570, 830, 820, _darken(acc2, 0.18), acc2, border_w=3)
        for lx in range(8):
            draw.line([(70 + lx, 570), (70 + lx, 820)],
                      fill=_blend(acc2, _darken(acc2, 0.3), lx / 8))
        _draw_text_wrapped(draw, node["yes"], 100, 600, 700, body_font, _WHITE)

        # NO outcome box
        _draw_card(draw, _W - 830, 570, _W - 70, 820, _darken(acc3, 0.18), acc3, border_w=3)
        for lx in range(8):
            draw.line([(_W - 830 + lx, 570), (_W - 830 + lx, 820)],
                      fill=_blend(acc3, _darken(acc3, 0.3), lx / 8))
        _draw_text_wrapped(draw, node["no"], _W - 800, 600, 700, body_font, _WHITE)

    elif t == "problem_solution":
        problem  = item.get("problem",  "")
        solution = item.get("solution", "")
        mid = _W // 2

        # PROBLEM card
        _draw_card(draw, 50, 175, mid - 50, 920, _darken(acc3, 0.18), acc3, border_w=3)
        for lx in range(10):
            draw.line([(50 + lx, 175), (50 + lx, 920)],
                      fill=_blend(acc3, _darken(acc3, 0.4), lx / 10))
        draw.text((90, 195), "⚠  PROBLEM", font=label_font, fill=acc3)
        draw.line([(90, 240), (mid - 90, 240)], fill=acc3, width=2)
        _draw_text_wrapped(draw, problem, 90, 265, mid - 180, body_font, _OFF_WHITE)

        # Arrow divider
        for ay in range(400, 680, 2):
            draw.line([(mid - 4, ay), (mid + 4, ay)],
                      fill=_blend(acc3, acc2, (ay - 400) / 280))
        draw.text((mid - 28, 510), "→", font=_font(80, bold=True), fill=_blend(acc3, acc2, 0.5))

        # SOLUTION card
        _draw_card(draw, mid + 50, 175, _W - 50, 920, _darken(acc2, 0.18), acc2, border_w=3)
        for lx in range(10):
            draw.line([(mid + 50 + lx, 175), (mid + 50 + lx, 920)],
                      fill=_blend(acc2, _darken(acc2, 0.4), lx / 10))
        draw.text((mid + 90, 195), "✓  SOLUTION", font=label_font, fill=acc2)
        draw.line([(mid + 90, 240), (_W - 90, 240)], fill=acc2, width=2)
        _draw_text_wrapped(draw, solution, mid + 90, 265, mid - 180, body_font, _OFF_WHITE)

    # ── Footer ────────────────────────────────────────────────────────────────
    for fx in range(_W):
        draw.line([(fx, _H - 50), (fx, _H - 48)],
                  fill=_blend(acc1, acc2, fx / _W))
    draw.text((40, _H - 40), "BRollAI  ·  AI-Generated Visual Explainer",
              font=small_font, fill=_MUTED)

    img.save(png_path, "PNG")


def pptx_to_png(pptx_path: str, png_path: str, item: dict = None):
    """Convert PPTX slide to PNG. Uses LibreOffice if available, else Pillow."""
    lo = shutil.which("libreoffice") or shutil.which("soffice")
    if lo:
        out_dir = os.path.dirname(png_path)
        result = subprocess.run(
            [lo, "--headless", "--convert-to", "png", "--outdir", out_dir, pptx_path],
            capture_output=True, text=True, timeout=60
        )
        base   = os.path.splitext(os.path.basename(pptx_path))[0]
        lo_out = os.path.join(out_dir, base + ".png")
        if os.path.exists(lo_out) and lo_out != png_path:
            os.rename(lo_out, png_path)
        if os.path.exists(png_path):
            return

    if item:
        render_broll_item_to_png(item, png_path)
        return

    raise RuntimeError("LibreOffice not found and no item data supplied for Pillow fallback.")


